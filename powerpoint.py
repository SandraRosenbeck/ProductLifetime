from pptx import Presentation
import pandas as pd
import os
from zipfile import ZipFile
from PIL import Image
import io

ITEM_TYPE_MAP = {
    "fixed": 1,
    "waiting": 2,
    "not_fixed": 3
}

def extract_mpo_images_from_pptx(pptx_path, group_id, name, item_type, image_dir, next_id):
    mpo_image_paths = []
    seen_hashes = set()
    item_type_number = ITEM_TYPE_MAP.get(item_type, 0)

    with ZipFile(pptx_path, 'r') as zipf:
        index = 1
        for zipinfo in zipf.infolist():
            if zipinfo.filename.startswith("ppt/media/"):
                image_data = zipf.read(zipinfo)
                img_hash = hash(image_data)
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)
                if b"MPF" in image_data[:64]:
                    try:
                        img = Image.open(io.BytesIO(image_data))
                        img.seek(0)

                        output_filename = f"group{group_id}_{next_id}_{name}_{item_type_number}.jpg"
                        output_path = os.path.join(image_dir, output_filename)
                        img.convert("RGB").save(output_path, "JPEG")
                        print(f"Extracted MPO as JPEG: {output_path}")

                        mpo_image_paths.append(output_path)
                        break  # Stop after first successful extraction
                    except Exception as e:
                        print(f"Failed to process MPO ({zipinfo.filename}): {e}")
            index += 1

    return mpo_image_paths

def process_slide(slide, group_id, name, item_type, image_dir, pptx_path=None, next_id=None):
    body_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body_text += shape.text.strip()

    image_paths = []

    for shape in slide.shapes:
        image = getattr(shape, "image", None)
        if image:
            try:
                ext = image.ext
            except ValueError as e:
                print(f"Unsupported image format in group {group_id}, slide '{item_type}' by {name}: {e}")

                if pptx_path:
                    # Attempting MPO extraction due to unsupported image format
                    mpo_images = extract_mpo_images_from_pptx(pptx_path, group_id, name, item_type, image_dir, next_id)
                    image_paths.extend(mpo_images)
                continue

            item_type_number = ITEM_TYPE_MAP.get(item_type, 0)
            filename = f"group{group_id}_{next_id}_{name}_{item_type_number}.{ext}"
            full_path = os.path.join(image_dir, filename)

            with open(full_path, "wb") as f:
                f.write(image.blob)

            image_paths.append(full_path)
            break

    return body_text, image_paths


def parse_person_info(person_info):
    data = {}
    for line in person_info.split('\n'):
        if ':' in line:
            key, value = line.split(':', 1)
            data[key.strip()] = value.strip()
    return data


def parse_pptx(path, group_id, image_dir, next_id):
    prs = Presentation(path)
    slide_list = list(prs.slides)

    data = []

    if len(slide_list) % 4 != 0:
        # deletes first slide (introduction) from those who have it
        del slide_list[0]

    for i in range(0, len(slide_list), 4):
        s1, s2, s3, s4 = slide_list[i:i+4]

        person_info = s1.shapes.placeholders[1].text.strip()
        info = parse_person_info(person_info)
        person_name = info.get('Navn', 'unknown').replace(" ", "_")

        fixed_body, fixed_image = process_slide(s2, group_id, person_name, "fixed", image_dir, path, next_id)
        waiting_body, waiting_image = process_slide(s3, group_id, person_name, "waiting", image_dir, path, next_id)
        notfixed_body, notfixed_image = process_slide(s4, group_id, person_name, "not_fixed", image_dir, path, next_id)

        row = {
            "id": next_id,
            "group": group_id,
            "person name": person_name,
            "age":  info.get('Alder'),
            "gender": info.get('Køn'),
            "postalcode":  info.get('Postnummer'),

            "fixed title": "En repareret ting",
            "fixed body": fixed_body,
            "fixed image dir": fixed_image[0] if fixed_image else None,

            "waiting title": "En ting, der venter på, eller er i gang med at blive repareret",
            "waiting body": waiting_body,
            "waiting image dir": waiting_image[0] if waiting_image else None,

            "not fixed title": "En ting, der ikke blev repareret/en ting der blev udskiftet",
            "not fixed body": notfixed_body,
            "not fixed image dir": notfixed_image[0] if notfixed_image else None,
        }
        data.append(row)
        next_id+=1
    return data


all_data = []
pptx_dir = "powerpointdata"
image_dir = "images"
os.makedirs(image_dir, exist_ok=True)
next_id=0

for filename in os.listdir(pptx_dir):
    if filename.endswith(".pptx"):
        group_id = filename[:2]
        path = os.path.join(pptx_dir, filename)

        if group_id in ("04", "06"):
            continue

        print('processing group', group_id)
        group_data = parse_pptx(path, group_id, image_dir, next_id)
        next_id += len(group_data)
        all_data.extend(group_data)

df = pd.DataFrame(all_data)
df.to_csv("powerpoint_data.csv", index=False)
